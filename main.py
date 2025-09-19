from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
from xml.dom import minidom
import argparse
import os
import json
from openai import OpenAI
from dotenv import load_dotenv
from pathlib import Path
import sys
from concurrent.futures import ThreadPoolExecutor
from typing import Dict, Optional, List, Tuple
import threading
import io
import base64
import tempfile
try:
    import pytesseract
    import cv2
    import numpy as np
    from PIL import Image, ImageDraw, ImageFont
    
    # Set Tesseract data path
    import os
    if not os.environ.get('TESSDATA_PREFIX'):
        os.environ['TESSDATA_PREFIX'] = '/usr/share/tessdata'
    
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("Warning: OCR libraries not installed. Image text extraction will be skipped.")
    print("Install with: pip install pytesseract Pillow opencv-python numpy")

# Load environment variables
load_dotenv()

# Initialize OpenAI client
client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY")
)

# Thread-safe translation cache
translation_cache: Dict[str, str] = {}
cache_lock = threading.Lock()

# OCR Text Detection and Extraction Functions
def extract_image_from_shape(shape) -> Optional[bytes]:
    """Extract image bytes from a PowerPoint shape."""
    try:
        if hasattr(shape, 'image'):
            return shape.image.blob
        elif hasattr(shape, 'part') and hasattr(shape.part, 'blob'):
            return shape.part.blob
        return None
    except Exception as e:
        print(f"Error extracting image from shape: {str(e)}")
        return None

def perform_ocr_on_image(image_bytes: bytes) -> List[Dict]:
    """Perform OCR on image and return text with bounding boxes."""
    if not OCR_AVAILABLE:
        return []
    
    try:
        # Convert bytes to PIL Image
        image = Image.open(io.BytesIO(image_bytes))
        
        # Handle different image formats
        if image.format in ['WMF', 'EMF']:
            print(f"Skipping {image.format} format image (not supported for OCR)")
            return []
        
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Check if image is large enough for OCR
        if image.width < 50 or image.height < 50:
            print(f"Image too small for OCR ({image.width}x{image.height})")
            return []
        
        # Convert PIL to OpenCV format
        cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        
        # Enhance image for better OCR (optional preprocessing)
        # Convert to grayscale for better text detection
        gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
        
        # Apply threshold to get better contrast
        _, threshold_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Use pytesseract to get detailed OCR data on both original and processed images
        try:
            # Try with original image first
            ocr_data = pytesseract.image_to_data(cv_image, output_type=pytesseract.Output.DICT, config='--psm 6')
        except:
            try:
                # Fallback to threshold image
                ocr_data = pytesseract.image_to_data(threshold_img, output_type=pytesseract.Output.DICT, config='--psm 6')
            except:
                # Last resort with different PSM
                ocr_data = pytesseract.image_to_data(cv_image, output_type=pytesseract.Output.DICT, config='--psm 3')
        
        # Extract text blocks with their positions
        text_blocks = []
        n_boxes = len(ocr_data['level'])
        
        for i in range(n_boxes):
            confidence = int(ocr_data['conf'][i])
            text = ocr_data['text'][i].strip()
            
            # Only include text with reasonable confidence and non-empty content
            if confidence > 30 and text and len(text) > 1:  # At least 2 characters
                text_block = {
                    'text': text,
                    'confidence': confidence,
                    'left': ocr_data['left'][i],
                    'top': ocr_data['top'][i],
                    'width': ocr_data['width'][i],
                    'height': ocr_data['height'][i],
                    'level': ocr_data['level'][i]
                }
                text_blocks.append(text_block)
        
        if text_blocks:
            print(f"Found {len(text_blocks)} text blocks in image ({image.width}x{image.height})")
            for block in text_blocks:
                print(f"  '{block['text']}' (confidence: {block['confidence']}%)")
        
        # Group text blocks by level to form coherent text regions
        grouped_blocks = group_text_blocks(text_blocks)
        return grouped_blocks
        
    except Exception as e:
        print(f"Error performing OCR: {str(e)}")
        return []

def group_text_blocks(text_blocks: List[Dict]) -> List[Dict]:
    """Group nearby text blocks into coherent regions."""
    if not text_blocks:
        return []
    
    # Sort blocks by their position (top to bottom, left to right)
    text_blocks.sort(key=lambda x: (x['top'], x['left']))
    
    grouped_blocks = []
    current_group = []
    
    for block in text_blocks:
        if not current_group:
            current_group = [block]
        else:
            # Check if this block should be grouped with the current group
            last_block = current_group[-1]
            
            # Group blocks that are on the same line or very close vertically
            vertical_distance = abs(block['top'] - last_block['top'])
            horizontal_gap = block['left'] - (last_block['left'] + last_block['width'])
            
            if vertical_distance < 20 and horizontal_gap < 50:  # Same line
                current_group.append(block)
            else:
                # Finalize current group and start new one
                if current_group:
                    grouped_blocks.append(merge_text_blocks(current_group))
                current_group = [block]
    
    # Don't forget the last group
    if current_group:
        grouped_blocks.append(merge_text_blocks(current_group))
    
    return grouped_blocks

def merge_text_blocks(blocks: List[Dict]) -> Dict:
    """Merge a list of text blocks into a single text region."""
    if not blocks:
        return {}
    
    if len(blocks) == 1:
        return blocks[0]
    
    # Calculate bounding box for all blocks
    min_left = min(block['left'] for block in blocks)
    min_top = min(block['top'] for block in blocks)
    max_right = max(block['left'] + block['width'] for block in blocks)
    max_bottom = max(block['top'] + block['height'] for block in blocks)
    
    # Combine text with spaces
    combined_text = ' '.join(block['text'] for block in blocks)
    
    # Calculate average confidence
    avg_confidence = sum(block['confidence'] for block in blocks) / len(blocks)
    
    return {
        'text': combined_text,
        'confidence': avg_confidence,
        'left': min_left,
        'top': min_top,
        'width': max_right - min_left,
        'height': max_bottom - min_top,
        'level': min(block['level'] for block in blocks)
    }

def overlay_text_on_image(image_bytes: bytes, text_regions: List[Dict]) -> bytes:
    """Overlay translated text back onto the image."""
    if not OCR_AVAILABLE or not text_regions:
        return image_bytes
    
    try:
        # Load original image
        image = Image.open(io.BytesIO(image_bytes))
        draw = ImageDraw.Draw(image)
        
        for region in text_regions:
            if 'translated_text' not in region:
                continue
                
            text = region['translated_text']
            x, y = region['left'], region['top']
            width, height = region['width'], region['height']
            
            # Try to find appropriate font size
            font_size = estimate_font_size(text, width, height)
            
            try:
                # Try to use a system font
                font = ImageFont.truetype("arial.ttf", font_size)
            except:
                try:
                    font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", font_size)
                except:
                    try:
                        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
                    except:
                        font = ImageFont.load_default()
            
            # Create a background rectangle to cover original text
            background_color = estimate_background_color(image, x, y, width, height)
            draw.rectangle([x, y, x + width, y + height], fill=background_color)
            
            # Draw the translated text
            text_color = (0, 0, 0)  # Black text by default
            draw.text((x + 2, y + 2), text, font=font, fill=text_color)
        
        # Convert back to bytes
        output_buffer = io.BytesIO()
        image.save(output_buffer, format='PNG')
        return output_buffer.getvalue()
        
    except Exception as e:
        print(f"Error overlaying text on image: {str(e)}")
        return image_bytes

def estimate_font_size(text: str, width: int, height: int) -> int:
    """Estimate appropriate font size for given text and region dimensions."""
    # Simple heuristic: assume average character width is 60% of font size
    if not text:
        return 12
    
    char_count = len(text)
    estimated_font_size = int(width * 0.8 / (char_count * 0.6))
    
    # Constrain font size to reasonable bounds
    font_size = max(8, min(estimated_font_size, height - 4, 72))
    return font_size

def estimate_background_color(image: Image.Image, x: int, y: int, width: int, height: int) -> Tuple[int, int, int]:
    """Estimate the background color of a text region."""
    try:
        # Sample pixels around the text region border
        region = image.crop((max(0, x-2), max(0, y-2), 
                           min(image.width, x + width + 2), 
                           min(image.height, y + height + 2)))
        
        # Convert to RGB if necessary
        if region.mode != 'RGB':
            region = region.convert('RGB')
        
        # Get the most common color (simple approach)
        colors = region.getcolors(maxcolors=256)
        if colors:
            # Return the most frequent color
            most_common_color = max(colors, key=lambda x: x[0])[1]
            return most_common_color
        
    except Exception as e:
        print(f"Error estimating background color: {str(e)}")
    
    # Default to white background
    return (255, 255, 255)

def translate_text(text: str, source_lang: str = 'zh', target_lang: str = 'en') -> str:
    """Translate text from source_lang to target_lang using OpenAI with rate limiting."""
    text = text.strip()
    if not text:
        return text
    
    # Check cache first
    with cache_lock:
        if text in translation_cache:
            return translation_cache[text]
    
    try:
        chunks = chunk_text(text)
        translated_chunks = []
        for chunk in chunks:
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    response = client.chat.completions.create(
                        model="gpt-3.5-turbo",
                        messages=[
                            {"role": "system", "content": f"You are a translator. Translate the following {source_lang} text to {target_lang}. Provide only the translated text without any additional comments, explanations, or formatting like asterisks or markdown. Maintain a natural, fluent translation.IMPORTANT: Preserve all numbers, numerical values, dates, times, formatting, line breaks, and special characters exactly as they appear in the original text."},
                            {"role": "user", "content": chunk}
                        ],
                        temperature=0.3,
                        stream=False
                    )
                    translated_chunks.append(response.choices[0].message.content.strip())
                    break  # Success, exit retry loop
                    
                except Exception as e:
                    if "rate_limit_exceeded" in str(e) and attempt < max_retries - 1:
                        import time
                        wait_time = 2 ** attempt  # Exponential backoff: 1s, 2s, 4s
                        print(f"Rate limit hit, waiting {wait_time}s before retry...")
                        time.sleep(wait_time)
                    else:
                        raise e  # Re-raise if not rate limit or max retries reached
        
        translated = ' '.join(translated_chunks)
        
        # If translated looks like fallback message, return original
        if "provide the" in translated.lower() or "no text" in translated.lower():
            return text
        
        # Cache the successful translation
        with cache_lock:
            translation_cache[text] = translated
        
        return translated
        
    except Exception as e:
        print(f"Translation error: {str(e)}")
        return text

def get_alignment_value(alignment_str):
    """Convert alignment string to PP_ALIGN enum value."""
    alignment_map = {
        'PP_ALIGN.CENTER': PP_ALIGN.CENTER,
        'PP_ALIGN.LEFT': PP_ALIGN.LEFT,
        'PP_ALIGN.RIGHT': PP_ALIGN.RIGHT,
        'PP_ALIGN.JUSTIFY': PP_ALIGN.JUSTIFY,
        'None': None
    }
    return alignment_map.get(alignment_str)

def get_shape_properties(shape):
    """Extract all properties from a shape with improved color detection."""
    shape_data = {
        'text': '',
        'font_size': None,
        'font_name': None,
        'alignment': None,
        'width': shape.width,
        'height': shape.height,
        'left': shape.left,
        'top': shape.top,
        'bold': None,
        'italic': None,
        'line_spacing': None,
        'space_before': None,
        'space_after': None,
        'font_color': None,
        'font_color_type': None,
        'bullet': None,
        'level': None,
        'original_text_length': 0,
        'runs_data': []  # Store individual run properties
    }
    
    if hasattr(shape, "text"):
        shape_data['text'] = shape.text.strip()
        shape_data['original_text_length'] = len(shape_data['text'])
        
        if hasattr(shape, 'text_frame'):
            # Process all paragraphs and runs for complete formatting
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                if paragraph.runs:
                    for run_idx, run in enumerate(paragraph.runs):
                        run_data = {
                            'text': run.text,
                            'font_size': None,
                            'font_name': None,
                            'bold': None,
                            'italic': None,
                            'font_color': None,
                            'font_color_type': None
                        }
                        
                        # Extract font properties
                        if hasattr(run.font, 'size') and run.font.size is not None:
                            run_data['font_size'] = run.font.size.pt
                            if para_idx == 0 and run_idx == 0:  # Use first run as default
                                shape_data['font_size'] = run.font.size.pt
                        
                        if hasattr(run.font, 'name') and run.font.name:
                            run_data['font_name'] = run.font.name
                            if para_idx == 0 and run_idx == 0:
                                shape_data['font_name'] = run.font.name
                        
                        if hasattr(run.font, 'bold'):
                            run_data['bold'] = run.font.bold
                            if para_idx == 0 and run_idx == 0:
                                shape_data['bold'] = run.font.bold
                        
                        if hasattr(run.font, 'italic'):
                            run_data['italic'] = run.font.italic
                            if para_idx == 0 and run_idx == 0:
                                shape_data['italic'] = run.font.italic
                        
                        # Improved color detection
                        if hasattr(run.font, 'color') and run.font.color is not None:
                            try:
                                # Try different color types
                                if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                    run_data['font_color'] = str(run.font.color.rgb)
                                    run_data['font_color_type'] = 'rgb'
                                elif hasattr(run.font.color, 'theme_color') and run.font.color.theme_color is not None:
                                    run_data['font_color'] = str(run.font.color.theme_color)
                                    run_data['font_color_type'] = 'theme'
                                elif hasattr(run.font.color, 'brightness') and run.font.color.brightness is not None:
                                    run_data['font_color_type'] = 'brightness'
                                    run_data['font_color'] = str(run.font.color.brightness)
                                
                                # Set shape default color from first run
                                if para_idx == 0 and run_idx == 0 and run_data['font_color']:
                                    shape_data['font_color'] = run_data['font_color']
                                    shape_data['font_color_type'] = run_data['font_color_type']
                            except Exception as e:
                                print(f"Warning: Could not extract color from run: {e}")
                        
                        shape_data['runs_data'].append(run_data)
                
                # Extract paragraph properties from first paragraph
                if para_idx == 0:
                    if hasattr(paragraph, 'line_spacing'):
                        shape_data['line_spacing'] = paragraph.line_spacing
                    if hasattr(paragraph, 'space_before'):
                        shape_data['space_before'] = paragraph.space_before
                    if hasattr(paragraph, 'space_after'):
                        shape_data['space_after'] = paragraph.space_after
                    if hasattr(paragraph, 'alignment'):
                        shape_data['alignment'] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
                    if hasattr(paragraph, 'bullet'):
                        shape_data['bullet'] = paragraph.bullet
                    if hasattr(paragraph, 'level'):
                        shape_data['level'] = paragraph.level

    return shape_data

def apply_shape_properties(shape, shape_data):
    """Apply saved properties to a shape with improved text fitting and color handling."""
    try:
        # Restore shape dimensions and position
        shape.width = shape_data['width']
        shape.height = shape_data['height']
        shape.left = shape_data['left']
        shape.top = shape_data['top']
        
        # Clear existing text
        shape.text = ""
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = shape_data['text']
        
        # Calculate appropriate font size for text fitting
        original_font_size = shape_data.get('font_size', 12)
        adjusted_font_size = calculate_fitted_font_size(
            shape_data['text'], 
            shape_data.get('original_text_length', len(shape_data['text'])),
            original_font_size,
            shape_data['width'],
            shape_data['height']
        )
        
        # Apply font properties
        if adjusted_font_size:
            run.font.size = Pt(adjusted_font_size)
        
        run.font.name = shape_data['font_name'] or 'Arial'
        
        # Improved color application
        if shape_data.get('font_color'):
            try:
                color_type = shape_data.get('font_color_type', 'rgb')
                if color_type == 'rgb':
                    run.font.color.rgb = RGBColor.from_string(shape_data['font_color'])
                elif color_type == 'theme':
                    # For theme colors, we'll fall back to RGB if possible
                    try:
                        run.font.color.theme_color = int(shape_data['font_color'])
                    except:
                        # Fallback to black if theme color fails
                        run.font.color.rgb = RGBColor(0, 0, 0)
                else:
                    # Default fallback
                    run.font.color.rgb = RGBColor.from_string(shape_data['font_color'])
            except Exception as e:
                print(f"Warning: Could not apply font color {shape_data['font_color']}: {e}")
                # Fallback to black text
                run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Apply other font properties
        if shape_data['bold'] is not None:
            run.font.bold = shape_data['bold']
        if shape_data['italic'] is not None:
            run.font.italic = shape_data['italic']
        
        # Apply paragraph properties
        if shape_data['alignment']:
            paragraph.alignment = get_alignment_value(shape_data['alignment'])
        if shape_data['line_spacing']:
            paragraph.line_spacing = shape_data['line_spacing']
        if shape_data['space_before']:
            paragraph.space_before = shape_data['space_before']
        if shape_data['space_after']:
            paragraph.space_after = shape_data['space_after']
        if shape_data.get('bullet') is not None:
            paragraph.bullet = shape_data['bullet']
        if shape_data.get('level') is not None:
            paragraph.level = shape_data['level']
            
    except Exception as e:
        print(f"Error applying shape properties: {str(e)}")

def calculate_fitted_font_size(translated_text: str, original_length: int, original_font_size: float, shape_width: int, shape_height: int) -> float:
    """Calculate appropriate font size to fit translated text within shape boundaries."""
    try:
        if not translated_text or not original_length or not original_font_size:
            return original_font_size
        
        # Calculate text length ratio
        length_ratio = len(translated_text) / original_length
        
        # If translated text is significantly longer, reduce font size
        if length_ratio > 1.2:  # 20% longer
            # Reduce font size proportionally, but not too much
            size_reduction_factor = min(0.8, 1.0 / (length_ratio * 0.8))
            adjusted_size = original_font_size * size_reduction_factor
            
            # Don't go below 8pt or above original size
            adjusted_size = max(8, min(adjusted_size, original_font_size))
            
            print(f"Adjusted font size from {original_font_size}pt to {adjusted_size:.1f}pt (text {length_ratio:.1f}x longer)")
            return adjusted_size
        elif length_ratio < 0.7:  # 30% shorter
            # Slightly increase font size for much shorter text
            size_increase_factor = min(1.2, 1.0 / (length_ratio * 1.1))
            adjusted_size = original_font_size * size_increase_factor
            
            # Don't exceed 150% of original size
            adjusted_size = min(adjusted_size, original_font_size * 1.5)
            
            print(f"Adjusted font size from {original_font_size}pt to {adjusted_size:.1f}pt (text {length_ratio:.1f}x shorter)")
            return adjusted_size
        
        # For similar lengths, keep original size
        return original_font_size
        
    except Exception as e:
        print(f"Error calculating fitted font size: {e}")
        return original_font_size

def get_image_properties(shape):
    """Extract image properties and perform OCR if applicable."""
    image_data = {
        'width': shape.width,
        'height': shape.height,
        'left': shape.left,
        'top': shape.top,
        'ocr_regions': [],
        'image_data': None
    }
    
    try:
        # Extract image bytes
        image_bytes = extract_image_from_shape(shape)
        if image_bytes:
            # Store image as base64 for XML serialization
            image_data['image_data'] = base64.b64encode(image_bytes).decode('utf-8')
            
            # Perform OCR to extract text regions
            if OCR_AVAILABLE:
                ocr_regions = perform_ocr_on_image(image_bytes)
                image_data['ocr_regions'] = ocr_regions
                print(f"Extracted {len(ocr_regions)} text regions from image")
            
    except Exception as e:
        print(f"Error processing image shape: {str(e)}")
    
    return image_data

def apply_image_properties(shape, image_data):
    """Apply OCR-translated text back to image."""
    try:
        # Restore image dimensions and position
        shape.width = image_data['width']
        shape.height = image_data['height']
        shape.left = image_data['left']
        shape.top = image_data['top']
        
        # If we have OCR regions with translations, overlay them on the image
        if image_data.get('ocr_regions') and image_data.get('image_data'):
            # Decode base64 image data
            image_bytes = base64.b64decode(image_data['image_data'])
            
            # Check if any region has translated text
            regions_with_translation = [r for r in image_data['ocr_regions'] if 'translated_text' in r]
            
            if regions_with_translation:
                # Overlay translated text on the image
                modified_image_bytes = overlay_text_on_image(image_bytes, regions_with_translation)
                
                # Save modified image to temporary file and replace in shape
                replace_image_in_shape(shape, modified_image_bytes)
                print(f"Applied translations to {len(regions_with_translation)} text regions in image")
        
    except Exception as e:
        print(f"Error applying image properties: {str(e)}")

def replace_image_in_shape(shape, new_image_bytes: bytes):
    """Replace the image in a PowerPoint shape with new image data."""
    try:
        # Create a temporary file for the new image
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
            temp_file.write(new_image_bytes)
            temp_path = temp_file.name
        
        try:
            # This is a workaround since python-pptx doesn't directly support replacing images
            # We need to access the underlying XML structure
            from pptx.parts.image import ImagePart
            from pptx.oxml import parse_xml
            
            # Get the slide and shape relationship
            slide_part = shape.part
            
            # Create new image part
            with open(temp_path, 'rb') as f:
                new_image_bytes = f.read()
            
            # Find existing image relationship and replace it
            if hasattr(shape, '_element') and hasattr(shape._element, 'blip_rId'):
                rel_id = shape._element.blip_rId
                if rel_id in slide_part.rels:
                    # Replace the image data in the existing relationship
                    image_part = slide_part.rels[rel_id].target_part
                    image_part._blob = new_image_bytes
                    print("Successfully replaced image in shape")
            
        finally:
            # Clean up temporary file
            os.unlink(temp_path)
            
    except Exception as e:
        print(f"Warning: Could not replace image in shape: {str(e)}")
        print("Image with translated text has been processed but may not appear in final PowerPoint")

def get_table_properties(table):
    """Extract complete table properties."""
    table_data = {
        'rows': len(table.rows),
        'cols': len(table.columns),
        'cells': []
    }
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                'text': cell.text.strip(),
                'original_text_length': len(cell.text.strip()),
                'font_size': None,
                'font_name': None,
                'alignment': None,
                'margin_left': cell.margin_left,
                'margin_right': cell.margin_right,
                'margin_top': cell.margin_top,
                'margin_bottom': cell.margin_bottom,
                'vertical_anchor': str(cell.vertical_anchor) if cell.vertical_anchor else None,
                'font_color': None,
                'font_color_type': None,
                'bullet': None,
                'level': None
            }
            if cell.text_frame.paragraphs:
                paragraph = cell.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if hasattr(run.font, 'size') and run.font.size is not None:
                        cell_data['font_size'] = run.font.size.pt
                    if hasattr(run.font, 'name'):
                        cell_data['font_name'] = run.font.name
                    if hasattr(run.font, 'bold'):
                        cell_data['bold'] = run.font.bold
                    if hasattr(run.font, 'italic'):
                        cell_data['italic'] = run.font.italic
                    
                    # Improved color detection for table cells
                    if hasattr(run.font, 'color') and run.font.color is not None:
                        try:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                cell_data['font_color'] = str(run.font.color.rgb)
                                cell_data['font_color_type'] = 'rgb'
                            elif hasattr(run.font.color, 'theme_color') and run.font.color.theme_color is not None:
                                cell_data['font_color'] = str(run.font.color.theme_color)
                                cell_data['font_color_type'] = 'theme'
                        except Exception as e:
                            print(f"Warning: Could not extract cell color: {e}")
                
                if hasattr(paragraph, 'alignment'):
                    cell_data['alignment'] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
                if hasattr(paragraph, 'bullet'):
                    cell_data['bullet'] = paragraph.bullet
                if hasattr(paragraph, 'level'):
                    cell_data['level'] = paragraph.level
            row_data.append(cell_data)
        table_data['cells'].append(row_data)
    return table_data

def apply_table_properties(table, table_data):
    """Apply saved table properties to a table."""
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                cell_data = table_data['cells'][row_idx][col_idx]
                cell.margin_left = cell_data['margin_left']
                cell.margin_right = cell_data['margin_right']
                cell.margin_top = cell_data['margin_top']
                cell.margin_bottom = cell_data['margin_bottom']
                if cell_data['vertical_anchor']:
                    cell.vertical_anchor = eval(cell_data['vertical_anchor'])
                cell.text = ""
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = cell_data['text']
                
                # Calculate fitted font size for table cells
                original_font_size = cell_data.get('font_size', 12)
                original_length = cell_data.get('original_text_length', len(cell_data['text']))
                
                # Estimate cell width/height for font fitting (approximate)
                cell_width = cell.width if hasattr(cell, 'width') else 100
                cell_height = cell.height if hasattr(cell, 'height') else 20
                
                adjusted_font_size = calculate_fitted_font_size(
                    cell_data['text'], 
                    original_length,
                    original_font_size,
                    cell_width,
                    cell_height
                )
                
                if adjusted_font_size:
                    run.font.size = Pt(adjusted_font_size)
                    
                run.font.name = cell_data['font_name'] or 'Arial'
                
                # Improved color handling for table cells
                if cell_data.get('font_color'):
                    try:
                        color_type = cell_data.get('font_color_type', 'rgb')
                        if color_type == 'rgb':
                            run.font.color.rgb = RGBColor.from_string(cell_data['font_color'])
                        elif color_type == 'theme':
                            try:
                                run.font.color.theme_color = int(cell_data['font_color'])
                            except:
                                run.font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            run.font.color.rgb = RGBColor.from_string(cell_data['font_color'])
                    except Exception as e:
                        print(f"Warning: Could not apply cell font color: {e}")
                        run.font.color.rgb = RGBColor(0, 0, 0)
                
                if 'bold' in cell_data:
                    run.font.bold = cell_data['bold']
                if 'italic' in cell_data:
                    run.font.italic = cell_data['italic']
                if cell_data['alignment']:
                    paragraph.alignment = get_alignment_value(cell_data['alignment'])
                if cell_data.get('bullet') is not None:
                    paragraph.bullet = cell_data['bullet']
                if cell_data.get('level') is not None:
                    paragraph.level = cell_data['level']
            except Exception as e:
                print(f"Error setting cell properties: {str(e)}")

def extract_shapes_recursively(shapes, slide_element, shape_path="", translate=False):
    """Recursively extract shapes including those within groups."""
    for shape_index, shape in enumerate(shapes):
        current_path = f"{shape_path}.{shape_index}" if shape_path else str(shape_index)
        
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_element = ET.SubElement(slide_element, "table_element")
            table_element.set("shape_index", current_path)
            table_data = get_table_properties(shape.table)
            if translate:
                for row in table_data['cells']:
                    for cell in row:
                        cell['text'] = translate_text(cell['text'])
            props_element = ET.SubElement(table_element, "properties")
            props_element.text = json.dumps(table_data, indent=2)
            
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Found PICTURE shape at path: {current_path}")
            image_element = ET.SubElement(slide_element, "image_element")
            image_element.set("shape_index", current_path)
            image_data = get_image_properties(shape)
            if translate and image_data.get('ocr_regions'):
                # Translate OCR-extracted text
                for region in image_data['ocr_regions']:
                    if region.get('text'):
                        region['translated_text'] = translate_text(region['text'])
                        print(f"Translated: '{region['text']}' -> '{region['translated_text']}'")
            props_element = ET.SubElement(image_element, "properties")
            props_element.text = json.dumps(image_data, indent=2)
            
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"Processing GROUP shape at path: {current_path} with {len(shape.shapes)} sub-shapes")
            # Recursively process shapes within the group
            extract_shapes_recursively(shape.shapes, slide_element, current_path, translate)
            
        elif hasattr(shape, "text"):
            text_element = ET.SubElement(slide_element, "text_element")
            text_element.set("shape_index", current_path)
            shape_data = get_shape_properties(shape)
            if translate:
                shape_data['text'] = translate_text(shape_data['text'])
            props_element = ET.SubElement(text_element, "properties")
            props_element.text = json.dumps(shape_data, indent=2)

def extract_text_from_slide(slide, slide_number, translate=False):
    """Extract all text elements from a slide."""
    slide_element = ET.Element("slide")
    slide_element.set("number", str(slide_number))
    
    # Use recursive extraction to handle nested shapes in groups
    extract_shapes_recursively(slide.shapes, slide_element, "", translate)
    
    return slide_element

def ppt_to_xml(ppt_path: str, translate: bool = False) -> Optional[str]:
    """Convert PowerPoint to XML with intermediate saves."""
    root = ET.Element("presentation")
    base_dir = Path(ppt_path).parent
    try:
        prs = Presentation(ppt_path)
        root.set("file_path", os.path.basename(ppt_path))
        with ThreadPoolExecutor(max_workers=4) as executor:
            future_to_slide = {
                executor.submit(extract_text_from_slide, slide, slide_number, translate): slide_number 
                for slide_number, slide in enumerate(prs.slides, 1)
            }
            for future in future_to_slide:
                slide_number = future_to_slide[future]
                try:
                    slide_element = future.result()
                    root.append(slide_element)
                    intermediate_path = base_dir / f"slide_{slide_number}_{'translated' if translate else 'original'}.xml"
                    xml_str = minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
                    with open(intermediate_path, 'w', encoding='utf-8') as f:
                        f.write(xml_str)
                except Exception as e:
                    print(f"Error processing slide {slide_number}: {str(e)}")
        return minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    except Exception as e:
        print(f"Error processing presentation: {str(e)}")
        return None

def apply_properties_recursively(shapes, xml_slide, shape_path=""):
    """Recursively apply properties to shapes including those within groups."""
    for shape_index, shape in enumerate(shapes):
        current_path = f"{shape_path}.{shape_index}" if shape_path else str(shape_index)
        
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_element = xml_slide.find(f".//table_element[@shape_index='{current_path}']")
            if table_element is not None:
                props_element = table_element.find("properties")
                if props_element is not None and props_element.text:
                    try:
                        table_data = json.loads(props_element.text)
                        apply_table_properties(shape.table, table_data)
                    except Exception as e:
                        print(f"Error applying table properties: {str(e)}")
                        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_element = xml_slide.find(f".//image_element[@shape_index='{current_path}']")
            if image_element is not None:
                props_element = image_element.find("properties")
                if props_element is not None and props_element.text:
                    try:
                        image_data = json.loads(props_element.text)
                        apply_image_properties(shape, image_data)
                        print(f"Applied image properties to shape at path: {current_path}")
                    except Exception as e:
                        print(f"Error applying image properties: {str(e)}")
                        
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Recursively process shapes within the group
            apply_properties_recursively(shape.shapes, xml_slide, current_path)
            
        elif hasattr(shape, "text"):
            text_element = xml_slide.find(f".//text_element[@shape_index='{current_path}']")
            if text_element is not None:
                props_element = text_element.find("properties")
                if props_element is not None and props_element.text:
                    try:
                        shape_data = json.loads(props_element.text)
                        apply_shape_properties(shape, shape_data)
                    except Exception as e:
                        print(f"Error applying shape properties: {str(e)}")

def create_translated_ppt(original_ppt_path, translated_xml_path, output_ppt_path):
    """Create a new PowerPoint with translated text."""
    try:
        prs = Presentation(original_ppt_path)
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        for slide_number, slide in enumerate(prs.slides, 1):
            xml_slide = root.find(f".//slide[@number='{slide_number}']")
            if xml_slide is None:
                continue
            # Use recursive processing to handle nested shapes in groups
            apply_properties_recursively(slide.shapes, xml_slide, "")
        prs.save(output_ppt_path)
        print(f"Translated PowerPoint saved to: {output_ppt_path}")
    except Exception as e:
        print(f"Error creating translated PowerPoint: {str(e)}")

def clean_path(path: str) -> str:
    """Remove quotes and handle escaped spaces in path."""
    path = path.strip("'\"")
    path = path.replace("\\ ", " ")
    path = path.replace("\\'", "'")
    return path

def translate_text_with_cache(text: str) -> str:
    """Translate text using a cache to avoid duplicate API calls."""
    if not text or text.isspace():
        return text
    with cache_lock:
        if text in translation_cache:
            return translation_cache[text]
    try:
        translated = translate_text(text)
        with cache_lock:
            translation_cache[text] = translated
        return translated
    except Exception as e:
        print(f"Translation error: {str(e)}")
        return text

def chunk_text(text: str, max_chunk_size: int = 1000) -> list[str]:
    """Split text into smaller chunks while preserving sentence boundaries."""
    if len(text) <= max_chunk_size:
        return [text]
    chunks = []
    current_chunk = []
    current_size = 0
    # naive sentence splitting
    sentences = text.replace('。', '.').replace('！', '!').replace('？', '?').split('.')
    for sentence in sentences:
        sentence = sentence.strip() + '.'
        if current_size + len(sentence) > max_chunk_size and current_chunk:
            chunks.append(''.join(current_chunk))
            current_chunk = []
            current_size = 0
        current_chunk.append(sentence)
        current_size += len(sentence)
    if current_chunk:
        chunks.append(''.join(current_chunk))
    return chunks

def cleanup_intermediate_files(base_dir: Path, pattern: str = "slide_*.xml"):
    """Clean up intermediate XML files after successful processing."""
    try:
        for file in base_dir.glob(pattern):
            file.unlink()
    except Exception as e:
        print(f"Warning: Could not clean up intermediate files: {str(e)}")

def process_ppt_file(ppt_path: Path, source_lang: str, target_lang: str):
    """Process a single PPT/PPTX file from XML extraction to final translation."""
    try:
        if not ppt_path.is_file():
            print(f"Error: '{ppt_path}' is not a valid file.")
            return
        if ppt_path.suffix.lower() not in ['.ppt', '.pptx']:
            print(f"Error: '{ppt_path}' is not a PowerPoint file.")
            return

        base_dir = ppt_path.parent

        # Adjust global translate_text to use chosen source/target
        global translate_text
        original_translate_func = translate_text
        translate_text = lambda text: original_translate_func(text, source_lang, target_lang)

        # Original XML
        print(f"Generating original XML for {ppt_path.name}...")
        original_xml = ppt_to_xml(str(ppt_path), translate=False)
        if original_xml:
            original_output_path = base_dir / f"{ppt_path.stem}_original.xml"
            with open(original_output_path, 'w', encoding='utf-8') as f:
                f.write(original_xml)
            print(f"Original XML saved: {original_output_path}")

        # Translated XML
        print(f"Generating translated XML (from {source_lang} to {target_lang}) for {ppt_path.name}...")
        translated_xml = ppt_to_xml(str(ppt_path), translate=True)
        if translated_xml:
            translated_output_path = base_dir / f"{ppt_path.stem}_translated.xml"
            with open(translated_output_path, 'w', encoding='utf-8') as f:
                f.write(translated_xml)
            print(f"Translated XML saved: {translated_output_path}")

            # Build final PPT
            print(f"Creating translated PPT for {ppt_path.name}...")
            output_filename = f"{ppt_path.stem}_translated{ppt_path.suffix}"
            output_ppt_path = base_dir / output_filename
            create_translated_ppt(str(ppt_path), str(translated_output_path), str(output_ppt_path))

            # Cleanup
            cleanup_intermediate_files(base_dir)
            print("Cleanup complete.")

        # Restore original translate_text
        translate_text = original_translate_func

    except Exception as e:
        print(f"Error in process_ppt_file for {ppt_path}: {str(e)}")

def main():
    try:
        print("=== PowerPoint Translator with OCR Support ===")
        print(f"OCR Available: {'Yes' if OCR_AVAILABLE else 'No (install pytesseract, Pillow, opencv-python, numpy)'}")
        print()
        
        path_input = input("Enter path to a PPTX file OR directory: ").strip()
        path_input = clean_path(path_input)
        path_input = os.path.expanduser(path_input)
        source_lang = input("Enter source language code (default 'zh'): ").strip().lower() or 'zh'
        target_lang = input("Enter target language code (default 'en'): ").strip().lower() or 'en'
        
        # OCR configuration
        if OCR_AVAILABLE:
            ocr_enabled = input("Enable OCR for images? (y/n, default 'y'): ").strip().lower()
            ocr_enabled = ocr_enabled != 'n'
            
            if ocr_enabled:
                print("\nOCR Settings:")
                print("- Text confidence threshold: 30%")
                print("- Will extract text from images and translate it")
                print("- Translated text will be overlaid back onto images")
        else:
            ocr_enabled = False
            print("\nWarning: OCR libraries not installed. Image text will not be processed.")
        
        target_path = Path(path_input).resolve()

        if target_path.is_dir():
            print(f"\nProcessing directory: {target_path}")
            # Recursively process all .ppt or .pptx
            for root, dirs, files in os.walk(target_path):
                for file in files:
                    if file.lower().endswith(('.ppt', '.pptx')):
                        full_path = Path(root) / file
                        process_ppt_file(full_path, source_lang, target_lang)
        else:
            print(f"\nProcessing file: {target_path}")
            # Process single file
            process_ppt_file(target_path, source_lang, target_lang)
        
        print("\n=== Translation Complete ===")

    except Exception as e:
        print(f"Fatal error: {str(e)}")
        sys.exit(1)

def test_ocr_functionality():
    """Test OCR functionality with a sample image."""
    if not OCR_AVAILABLE:
        print("OCR libraries not available for testing")
        return False
    
    try:
        # Create a simple test image with text
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a white image with black text
        img = Image.new('RGB', (300, 100), color='white')
        draw = ImageDraw.Draw(img)
        
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()
        
        draw.text((10, 10), "Hello World", fill='black', font=font)
        draw.text((10, 40), "Test OCR", fill='black', font=font)
        
        # Convert to bytes
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_bytes = img_buffer.getvalue()
        
        # Test OCR
        ocr_results = perform_ocr_on_image(img_bytes)
        
        print("OCR Test Results:")
        for i, region in enumerate(ocr_results):
            print(f"  Region {i+1}: '{region['text']}' (confidence: {region['confidence']:.1f}%)")
        
        # Test translation overlay
        if ocr_results:
            # Add mock translations
            for region in ocr_results:
                region['translated_text'] = f"[TRANSLATED] {region['text']}"
            
            # Test overlay
            modified_bytes = overlay_text_on_image(img_bytes, ocr_results)
            print(f"Successfully overlaid {len(ocr_results)} text regions")
            
            return True
        else:
            print("No text detected in test image")
            return False
            
    except Exception as e:
        print(f"OCR test failed: {str(e)}")
        return False

if __name__ == "__main__":
    # Uncomment the line below to run OCR tests
    # test_ocr_functionality()
    main()
